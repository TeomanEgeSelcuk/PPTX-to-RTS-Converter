import os
import shutil
from pptx import Presentation
from pyth.plugins.rtf15.reader import Rtf15Reader
from pyth.plugins.plaintext.writer import PlaintextWriter

def pptx_to_rtf(pptx_file, rtf_file):
    presentation = Presentation(pptx_file)

    # Extract the text from the PowerPoint file
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text += shape.text_frame.text + "\n"

    # Write the extracted text to the RTF file
    with open(rtf_file, "w", encoding="utf-8") as rtf_output:
        doc = Rtf15Reader.read(text)
        rtf_output.write(PlaintextWriter.write(doc).getvalue())

def convert_pptx_to_rtf():
    # Create the 'converted' directory if it doesn't exist
    if not os.path.exists("converted"):
        os.makedirs("converted")

    # Get a list of all .pptx files in the current directory
    pptx_files = [file for file in os.listdir(".") if file.endswith(".pptx")]

    for pptx_file in pptx_files:
        rtf_file = os.path.join("converted", os.path.splitext(pptx_file)[0] + ".rtf")

        # Check if the corresponding .rtf file already exists
        if not os.path.exists(rtf_file):
            print(f"Converting {pptx_file} to {rtf_file}...")
            pptx_to_rtf(pptx_file, rtf_file)
        else:
            print(f"{rtf_file} already exists. Skipping conversion.")

if __name__ == "__main__":
    convert_pptx_to_rtf()
